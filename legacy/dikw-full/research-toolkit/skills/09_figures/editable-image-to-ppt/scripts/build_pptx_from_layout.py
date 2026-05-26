#!/usr/bin/env python
"""Build an editable PPTX from a slide layout JSON file."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:
    raise SystemExit(
        "Missing python-pptx. Install it in the active Python environment with: "
        "python -m pip install python-pptx"
    ) from exc


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build editable PPTX from layout JSON.")
    parser.add_argument("--layout", required=True, help="Layout JSON path.")
    parser.add_argument("--out", required=True, help="Output .pptx path.")
    parser.add_argument("--assets-root", help="Base directory for relative image paths.")
    return parser.parse_args()


def rgb(value: Any, default: RGBColor | None = None) -> RGBColor | None:
    if value is None:
        return default
    if isinstance(value, str):
        text = value.strip()
        if text.lower() in {"none", "transparent", "background"}:
            return None
        text = text.lstrip("#")
        if len(text) == 6:
            return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        return RGBColor(int(value[0]), int(value[1]), int(value[2]))
    raise ValueError(f"Unsupported color: {value!r}")


def align(value: str | None) -> PP_ALIGN:
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get((value or "center").lower(), PP_ALIGN.CENTER)


def valign(value: str | None) -> MSO_ANCHOR:
    return {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get((value or "middle").lower(), MSO_ANCHOR.MIDDLE)


def shape_type(value: str | None) -> MSO_SHAPE:
    return {
        "rect": MSO_SHAPE.RECTANGLE,
        "rectangle": MSO_SHAPE.RECTANGLE,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "ellipse": MSO_SHAPE.OVAL,
        "diamond": MSO_SHAPE.DIAMOND,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "trapezoid": MSO_SHAPE.TRAPEZOID,
    }.get((value or "rect").lower(), MSO_SHAPE.RECTANGLE)


def dash_style(value: str | None):
    if not value:
        return None
    return {
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "dashed": MSO_LINE_DASH_STYLE.DASH,
        "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "dotted": MSO_LINE_DASH_STYLE.ROUND_DOT,
    }.get(value.lower())


class Builder:
    def __init__(self, layout: dict[str, Any], out: Path, assets_root: Path):
        self.layout = layout
        self.out = out
        self.assets_root = assets_root
        self.px_w = float(layout.get("source_width") or layout.get("canvas", {}).get("width") or 1182)
        self.px_h = float(layout.get("source_height") or layout.get("canvas", {}).get("height") or 665)
        slide_size = layout.get("slide_size", {})
        self.slide_w_in = float(slide_size.get("width_in", 13.333333))
        self.slide_h_in = float(slide_size.get("height_in", 7.5))
        self.sx = self.slide_w_in / self.px_w
        self.sy = self.slide_h_in / self.px_h
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_w_in)
        self.prs.slide_height = Inches(self.slide_h_in)

    def x(self, value: float):
        return Inches(float(value) * self.sx)

    def y(self, value: float):
        return Inches(float(value) * self.sy)

    def w(self, value: float):
        return Inches(float(value) * self.sx)

    def h(self, value: float):
        return Inches(float(value) * self.sy)

    def set_font(self, run, spec: dict[str, Any]) -> None:
        font = spec.get("font") or spec.get("font_name") or "Microsoft YaHei"
        run.font.name = font
        run.font.size = Pt(float(spec.get("size", spec.get("font_size", 18))))
        run.font.bold = bool(spec.get("bold", False))
        run.font.italic = bool(spec.get("italic", False))
        col = rgb(spec.get("color", "#111111"))
        if col is not None:
            run.font.color.rgb = col
        r_pr = run._r.get_or_add_rPr()
        latin = r_pr.find(qn("a:latin"))
        if latin is None:
            latin = OxmlElement("a:latin")
            r_pr.append(latin)
        latin.set("typeface", font)
        ea = r_pr.find(qn("a:ea"))
        if ea is None:
            ea = OxmlElement("a:ea")
            r_pr.append(ea)
        ea.set("typeface", font)

    def image_path(self, value: str) -> Path:
        path = Path(value)
        if not path.is_absolute():
            path = self.assets_root / path
        return path

    def add_text(self, slide, el: dict[str, Any]) -> None:
        left, top, width, height = el["box"]
        shape = slide.shapes.add_textbox(self.x(left), self.y(top), self.w(width), self.h(height))
        shape.name = el.get("name", "text")
        tf = shape.text_frame
        tf.clear()
        margin = float(el.get("margin", 0))
        tf.margin_left = self.w(margin)
        tf.margin_right = self.w(margin)
        tf.margin_top = self.h(margin)
        tf.margin_bottom = self.h(margin)
        tf.vertical_anchor = valign(el.get("valign"))
        lines = str(el.get("text", "")).split("\n")
        for idx, line in enumerate(lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = align(el.get("align"))
            p.line_spacing = float(el.get("line_spacing", 1.05))
            run = p.add_run()
            run.text = line
            self.set_font(run, el)

    def add_image(self, slide, el: dict[str, Any]) -> None:
        left, top, width, height = el["box"]
        pic = slide.shapes.add_picture(
            str(self.image_path(el["path"])),
            self.x(left),
            self.y(top),
            width=self.w(width),
            height=self.h(height),
        )
        pic.name = el.get("name", Path(el["path"]).stem)
        if "rotation" in el:
            pic.rotation = float(el["rotation"])

    def add_shape(self, slide, el: dict[str, Any]) -> None:
        left, top, width, height = el["box"]
        shape = slide.shapes.add_shape(
            shape_type(el.get("shape")),
            self.x(left),
            self.y(top),
            self.w(width),
            self.h(height),
        )
        shape.name = el.get("name", el.get("shape", "shape"))
        fill = rgb(el.get("fill"))
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        line = rgb(el.get("line"), RGBColor(0, 0, 0))
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line
            shape.line.width = Pt(float(el.get("line_width", 1)))
        if len(shape.adjustments) and "radius" in el:
            shape.adjustments[0] = float(el["radius"])
        if "rotation" in el:
            shape.rotation = float(el["rotation"])
        if "shadow" in el:
            shape.shadow.inherit = bool(el["shadow"])
        elif el.get("shadow_off", True):
            shape.shadow.inherit = False

    def add_line(self, slide, el: dict[str, Any]) -> None:
        if "points" in el:
            x1, y1, x2, y2 = el["points"]
        else:
            left, top, width, height = el["box"]
            x1, y1, x2, y2 = left, top, left + width, top + height
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, self.x(x1), self.y(y1), self.x(x2), self.y(y2))
        line.name = el.get("name", "line")
        col = rgb(el.get("line", el.get("color", "#000000")))
        if col is not None:
            line.line.color.rgb = col
        line.line.width = Pt(float(el.get("line_width", el.get("weight", 1))))
        dash = dash_style(el.get("dash"))
        if dash:
            line.line.dash_style = dash

    def add_slide(self, spec: dict[str, Any]) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = rgb(spec.get("background", self.layout.get("background", "#FFFFFF")))
        if bg is not None:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg
        for el in spec.get("elements", []):
            kind = el.get("type", "shape").lower()
            if kind == "text":
                self.add_text(slide, el)
            elif kind == "image":
                self.add_image(slide, el)
            elif kind == "shape":
                self.add_shape(slide, el)
            elif kind == "line":
                self.add_line(slide, el)
            else:
                raise ValueError(f"Unsupported element type: {kind}")

    def build(self) -> None:
        slides = self.layout.get("slides")
        if not slides:
            slides = [{"elements": self.layout.get("elements", [])}]
        for spec in slides:
            self.add_slide(spec)
        self.out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(self.out)


def main() -> None:
    args = parse_args()
    layout_path = Path(args.layout)
    layout = json.loads(layout_path.read_text(encoding="utf-8-sig"))
    assets_root = Path(args.assets_root) if args.assets_root else layout_path.parent
    Builder(layout, Path(args.out), assets_root).build()
    print(args.out)


if __name__ == "__main__":
    main()
