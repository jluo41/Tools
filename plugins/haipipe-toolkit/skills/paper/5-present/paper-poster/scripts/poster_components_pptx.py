import fitz, os, tempfile, shutil
from pptx import Presentation
from pptx.util import Mm
from pptx.dml.color import RGBColor

doc = fitz.open('poster/main.pdf')
page = doc[0]
pw, ph = page.rect.width, page.rect.height

# A0 dimensions in mm (adjust for portrait/A1)
W_mm, H_mm = 1189, 841  # landscape
# W_mm, H_mm = 841, 1189  # portrait

def pts_to_mm(x, y):
    return x / pw * W_mm, y / ph * H_mm

# ── Define regions from tcbposter grid ──
# Format: name → (col_0based, row_start, col_span, row_end)
# rows=20, columns=4 for landscape (3 for portrait)
COLS = 4
row_h = ph / 20
col_w = pw / COLS

regions = {
    "title":        (0, 0, COLS, 4),
    "stats":        (0, 4, COLS, 6),
    # ... add one entry per posterbox, matching between=rowN and rowM
    # Example for 4-column landscape:
    "background":   (0, 6, 1, 11),
    "contributions":(0, 11, 1, 16),
    "references":   (0, 16, 1, 20),
    "paradigms":    (1, 6, 1, 11),
    "models":       (1, 11, 1, 20),
    "architecture": (2, 6, 1, 10),
    "results1":     (2, 10, 1, 20),
    "hallucination":(3, 6, 1, 11),
    "ablation":     (3, 11, 1, 15),
    "takeaways":    (3, 15, 1, 20),
}

# ── Create PPTX ──
prs = Presentation()
prs.slide_width = Mm(W_mm)
prs.slide_height = Mm(H_mm)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Set background
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF3, 0xFF)  # venue bg color

tmpdir = tempfile.mkdtemp()
mat = fitz.Matrix(300/72, 300/72)  # 300 DPI

for name, (col, r0, span, r1) in regions.items():
    # Clip rectangle in PDF points
    clip = fitz.Rect(col * col_w, r0 * row_h,
                     (col + span) * col_w, r1 * row_h)
    pix = page.get_pixmap(matrix=mat, clip=clip)
    img_path = os.path.join(tmpdir, f"{name}.png")
    pix.save(img_path)

    # Position in mm
    left, top = pts_to_mm(clip.x0, clip.y0)
    right, bottom = pts_to_mm(clip.x1, clip.y1)

    slide.shapes.add_picture(img_path, Mm(left), Mm(top),
                             Mm(right - left), Mm(bottom - top))

prs.save('poster/poster_components.pptx')
doc.close()
shutil.rmtree(tmpdir)
