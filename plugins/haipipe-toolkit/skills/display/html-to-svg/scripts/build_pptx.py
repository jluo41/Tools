#!/usr/bin/env python3
"""Assemble the NN-*.svg slides into a 16:9 .pptx — true vector SVG embedding.

Each slide gets the SVG as a native vector graphic (PowerPoint 2019/365)
with a high-res PNG fallback (older PowerPoint), plus speaker notes.

python-pptx cannot insert SVG directly, so we do the documented dual-image
trick: insert the PNG, then patch the picture's <a:blip> with the
svgBlip extension pointing at an SVG part added to the package.

Requires:  pip install python-pptx cairosvg
Run:       python3 build_pptx.py          → reach-adhd-po-briefing.pptx
"""
import glob
import os
from pathlib import Path

import cairosvg
from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.util import Inches

HERE = Path(__file__).parent
OUT_PPTX = HERE / "deck.pptx"
PNG_DIR = HERE / "_png_fallback"
PNG_WIDTH = 1920  # fallback raster resolution

NS_ASVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

# Speaker notes per slide, keyed by svg file stem (edit per deck, or leave empty)
NOTES = {
    # "01-cover": "One-sentence guidance for the presenter…",
}


def add_svg_picture(slide, prs, svg_path, png_path, next_media_idx):
    """Insert PNG full-bleed, then patch in the vector SVG via svgBlip."""
    pic = slide.shapes.add_picture(
        str(png_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    # add the SVG bytes as a package part and relate it to the slide
    partname = PackURI(f"/ppt/media/image_svg{next_media_idx}.svg")
    svg_part = Part(partname, "image/svg+xml",
                    slide.part.package, blob=svg_path.read_bytes())
    rid = slide.part.relate_to(svg_part, RT.IMAGE)

    # patch <a:blip>: append extLst/ext/asvg:svgBlip pointing at the SVG rel
    blip = pic._element.blipFill.blip
    ext_lst = blip.makeelement(f"{{{NS_A}}}extLst", {})
    ext = etree.SubElement(ext_lst, f"{{{NS_A}}}ext", {"uri": SVG_EXT_URI})
    etree.SubElement(ext, f"{{{NS_ASVG}}}svgBlip",
                     {f"{{{NS_R}}}embed": rid},
                     nsmap={"asvg": NS_ASVG})
    blip.append(ext_lst)


def main():
    import sys
    global OUT_PPTX
    src_dir = HERE
    if "--editable" in sys.argv:
        # one-line variant (generate_svgs.py --ppt): unwrapped sentences,
        # meant for Convert to Shape + manual re-wrap (figure-to-svg Lesson 16)
        src_dir = HERE / "ppt-editable"
        OUT_PPTX = HERE / (OUT_PPTX.stem + "-editable.pptx")
    PNG_DIR.mkdir(exist_ok=True)
    svgs = sorted(src_dir.glob("[0-9][0-9]-*.svg"))
    if not svgs:
        raise SystemExit(f"no NN-*.svg files found in {src_dir}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # fully blank layout

    for i, svg in enumerate(svgs, 1):
        png = PNG_DIR / (svg.stem + ".png")
        cairosvg.svg2png(url=str(svg), write_to=str(png), output_width=PNG_WIDTH)
        slide = prs.slides.add_slide(blank)
        add_svg_picture(slide, prs, svg, png, i)
        note = NOTES.get(svg.stem)
        if note:
            slide.notes_slide.notes_text_frame.text = note
        print(f"slide {i:2d}  {svg.name}")

    prs.save(OUT_PPTX)
    print(f"\nwrote {OUT_PPTX.name} ({OUT_PPTX.stat().st_size // 1024} KB, "
          f"{len(svgs)} slides, vector SVG + PNG fallback, speaker notes)")


if __name__ == "__main__":
    main()
