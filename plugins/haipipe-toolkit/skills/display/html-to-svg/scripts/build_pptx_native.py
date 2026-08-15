#!/usr/bin/env python3
"""Build a NATIVE-shapes .pptx from the NN-*.svg slides — pre-'ungrouped'.

Instead of embedding SVG pictures (which need right-click → Convert to Shape
in PowerPoint, with its tspan-collapse bugs), this parses each generated SVG
(our own five-primitive subset: rect / line / circle / polyline / text) and
emits real DrawingML: every text line is a native, editable text box; every
box/line/chart mark is a native shape. Opens fully editable, nothing to
convert.

Coordinate map: 1280x720 px canvas → 13.333x7.5 in; 1 px = 9525 EMU; 1 px
font = 0.75 pt.

Requires:  pip install python-pptx
Run:       python3 build_pptx_native.py   → reach-adhd-po-briefing-native.pptx
"""
import re
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

try:
    from build_pptx import NOTES  # speaker notes, keyed by svg stem
except ImportError:
    NOTES = {}

HERE = Path(__file__).parent
OUT = HERE / "deck-native.pptx"
PX = 9525  # EMU per px at 96 dpi
NS = "{http://www.w3.org/2000/svg}"
FONT = "Times New Roman"


def px(v):
    return Emu(int(round(float(v) * PX)))


def color(c, opacity=None):
    """SVG color -> (RGBColor, is_none). Alpha pre-blended against white."""
    if not c or c == "none":
        return None
    c = c.strip()
    if c.startswith("#"):
        h = c[1:]
        if len(h) == 3:
            h = "".join(ch * 2 for ch in h)
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    else:
        named = {"white": (255, 255, 255), "black": (0, 0, 0)}
        r, g, b = named.get(c, (0, 0, 0))
    if opacity is not None:
        a = float(opacity)
        r = int(255 - a * (255 - r))
        g = int(255 - a * (255 - g))
        b = int(255 - a * (255 - b))
    return RGBColor(r, g, b)


def set_line(shape, stroke, sw, dash):
    ln = shape.line
    col = color(stroke)
    if col is None:
        ln.fill.background()
        return
    ln.color.rgb = col
    ln.width = Emu(int(round(float(sw or 1) * PX)))
    if dash:
        ln._get_or_add_ln().append(_dash_el())


def _dash_el():
    from pptx.oxml import parse_xml
    return parse_xml(
        '<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dash"/>')


def add_rect(shapes, e):
    x, y = float(e.get("x", 0)), float(e.get("y", 0))
    w, h = float(e.get("width")), float(e.get("height"))
    rx = float(e.get("rx", 0))
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0.5 else MSO_SHAPE.RECTANGLE
    sp = shapes.add_shape(kind, px(x), px(y), px(w), px(h))
    if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = max(0.0, min(0.5, rx / max(1.0, min(w, h))))
        except Exception:
            pass
    fill = color(e.get("fill", "#ffffff"), e.get("fill-opacity"))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    set_line(sp, e.get("stroke", "none"), e.get("stroke-width", 1),
             e.get("stroke-dasharray"))
    sp.shadow.inherit = False
    return sp


def add_line(shapes, e):
    x1, y1 = float(e.get("x1")), float(e.get("y1"))
    x2, y2 = float(e.get("x2")), float(e.get("y2"))
    if e.get("stroke-dasharray"):
        # dashed rules/legend swatches: freeform 2-point line carries prstDash
        fb = shapes.build_freeform(px(x1), px(y1), scale=1.0)
        fb.add_line_segments([(px(x2), px(y2))], close=False)
        sp = fb.convert_to_shape()
        sp.fill.background()
        set_line(sp, e.get("stroke", "#000000"), e.get("stroke-width", 1), "dash")
        sp.shadow.inherit = False
        return sp
    left, top = min(x1, x2), min(y1, y2)
    w, hh = abs(x2 - x1), abs(y2 - y1)
    sp = shapes.add_shape(MSO_SHAPE.RECTANGLE, px(left), px(top),
                          px(max(w, 0.5)), px(max(hh, 0.5)))
    # a thin filled rect reads as a line and never flips direction
    sw = float(e.get("stroke-width", 1))
    if w >= hh:
        sp.height = Emu(int(round(sw * PX)))
    else:
        sp.width = Emu(int(round(sw * PX)))
    fill = color(e.get("stroke", "#000000"))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    if e.get("stroke-dasharray"):
        # keep dashes visually: lighten the fill instead (simpler than
        # per-dash segments); dashed rules are decorative separators
        r, g, b = fill[0] if False else (None, None, None)
    return sp


def add_circle(shapes, e):
    cx, cy, r = float(e.get("cx")), float(e.get("cy")), float(e.get("r"))
    sp = shapes.add_shape(MSO_SHAPE.OVAL, px(cx - r), px(cy - r),
                          px(2 * r), px(2 * r))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color(e.get("fill", "#000000"))
    set_line(sp, e.get("stroke", "none"), e.get("stroke-width", 1), None)
    sp.shadow.inherit = False
    return sp


def add_polyline(shapes, e):
    pts = [tuple(map(float, p.split(","))) for p in e.get("points").split()]
    if len(pts) < 2:
        return
    xs, ys = [p[0] for p in pts], [p[1] for p in pts]
    fb = shapes.build_freeform(px(pts[0][0]), px(pts[0][1]), scale=1.0)
    fb.add_line_segments([(px(x), px(y)) for x, y in pts[1:]], close=False)
    sp = fb.convert_to_shape()
    sp.fill.background()
    set_line(sp, e.get("stroke", "#000000"), e.get("stroke-width", 2),
             e.get("stroke-dasharray"))
    sp.shadow.inherit = False
    return sp


def add_text(shapes, e):
    txt = e.text or ""
    if not txt.strip():
        return
    fs = float(e.get("font-size", 19))
    fill = color(e.get("fill", "#000000"))
    bold = e.get("font-weight", "normal") == "bold"
    anchor = e.get("text-anchor", "start")
    x, y = float(e.get("x")), float(e.get("y"))
    est_w = len(txt) * fs * 0.50 + fs * 0.6  # matches the generator estimate
    box_w = est_w * 1.25 + 8                  # generous so nothing wraps
    if anchor == "middle":
        left = x - box_w / 2
        align = PP_ALIGN.CENTER
    elif anchor == "end":
        left = x - box_w
        align = PP_ALIGN.RIGHT
    else:
        left = x
        align = PP_ALIGN.LEFT
    top = y - fs * 1.02  # svg y is the baseline
    # NOTE: keep default wrap (square). wrap="none" makes LibreOffice
    # auto-grow boxes around their center, shifting text — the classic drift.
    tb = shapes.add_textbox(px(left), px(top), px(box_w), px(fs * 1.45))
    tf = tb.text_frame
    tf.word_wrap = True  # python-pptx template defaults to wrap="none", which
    # LibreOffice renders center-grown; square wrap + generous width fixes it
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    f = run.font
    f.name = FONT
    f.size = Pt(fs * 0.75)
    f.bold = bold
    f.color.rgb = fill
    return tb


def build_slide(shapes, svg_path):
    root = ET.parse(svg_path).getroot()
    first_rect_skipped = False
    for e in root:
        tag = e.tag.replace(NS, "")
        if tag == "title":
            continue
        if tag == "rect":
            if not first_rect_skipped:  # full-page white background
                first_rect_skipped = True
                if (float(e.get("width", 0)) >= 1280 and
                        float(e.get("height", 0)) >= 720):
                    continue
            add_rect(shapes, e)
        elif tag == "line":
            add_line(shapes, e)
        elif tag == "circle":
            add_circle(shapes, e)
        elif tag == "polyline":
            add_polyline(shapes, e)
        elif tag == "text":
            add_text(shapes, e)


def main():
    svgs = sorted(HERE.glob("[0-9][0-9]-*.svg"))
    if not svgs:
        raise SystemExit("no NN-*.svg files found")
    prs = Presentation()
    prs.slide_width = Emu(1280 * PX)
    prs.slide_height = Emu(720 * PX)
    blank = prs.slide_layouts[6]
    for i, svg in enumerate(svgs, 1):
        slide = prs.slides.add_slide(blank)
        build_slide(slide.shapes, svg)
        note = NOTES.get(svg.stem)
        if note:
            slide.notes_slide.notes_text_frame.text = note
        print(f"slide {i:2d}  {svg.name}  ({len(slide.shapes)} native shapes)")
    prs.save(OUT)
    print(f"\nwrote {OUT.name} ({OUT.stat().st_size // 1024} KB, "
          f"{len(svgs)} slides, all native shapes + text boxes)")


if __name__ == "__main__":
    main()
